#!/usr/bin/env python3
"""Bring docs/Capstone.pptx up to the code that was actually built.

The deck was written before the platform existed, so it describes intent: a
STIX 2.1 exporter, a phoneinfoga module, a four-week roadmap. Rather than
restyle it, this rewrites the wording of the existing slides against the
current source and clones its own layouts for the capabilities that were added
afterwards -- evidence preservation, per-tool parsers, redaction, resilience.

Cloning keeps every slide's fonts, colours and geometry identical to the
original design, so the deck stays one piece of visual work.

    python3 scripts/build_capstone_deck.py docs/Capstone.pptx
"""

from __future__ import annotations

import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TARGET_MODE as RTM
from pptx.opc.package import _Relationship
from pptx.util import Pt

# ---------------------------------------------------------------------------
# Text handling
# ---------------------------------------------------------------------------


def set_text(shape, text: str) -> None:
    """Replace a shape's text, keeping the first run's character formatting.

    Assigning to ``text_frame.text`` discards the run properties the design
    depends on (weight, colour, size), so the first run is reused and the rest
    dropped. Newlines become separate paragraphs cloned from the first, which
    is how the original multi-line boxes are built.
    """
    tf = shape.text_frame
    lines = text.split("\n")

    first = tf.paragraphs[0]
    if not first.runs:
        first.add_run()
    template_run = first.runs[0]

    # Drop every paragraph after the first, and every run after the first.
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    for run in list(first.runs[1:]):
        run._r.getparent().remove(run._r)

    template_run.text = lines[0]
    for line in lines[1:]:
        new_p = copy.deepcopy(first._p)
        first._p.getparent().append(new_p)
        para = tf.paragraphs[-1]
        for run in list(para.runs[1:]):
            run._r.getparent().remove(run._r)
        para.runs[0].text = line


def shape_by_text(slide, needle: str):
    """The first shape whose text contains ``needle``."""
    for shape in slide.shapes:
        if shape.has_text_frame and needle in shape.text_frame.text:
            return shape
    return None


def edit(slide, needle: str, text: str) -> bool:
    shape = shape_by_text(slide, needle)
    if shape is None:
        print(f"  !! no shape containing {needle!r}")
        return False
    set_text(shape, text)
    return True


def shrink(slide, needle: str, size: float) -> None:
    """Hold a rewritten block inside the box the design drew for it."""
    shape = shape_by_text(slide, needle)
    if shape is None:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)


# ---------------------------------------------------------------------------
# Slide cloning
# ---------------------------------------------------------------------------


def clone_slide(prs: Presentation, source_index: int):
    """Append a copy of an existing slide, images and all.

    python-pptx has no slide copy: the shape tree is deep-copied and every
    relationship the copy points at (the footer logo here) is re-created on the
    new slide, otherwise the picture references a part it cannot reach.
    """
    source = prs.slides[source_index]
    dest = prs.slides.add_slide(source.slide_layout)

    for shape in list(dest.shapes):
        shape._element.getparent().remove(shape._element)

    for shape in source.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shape._element))

    # The copied shapes carry the source's r:embed ids, so those exact ids have
    # to resolve on the new part -- hence writing the rel map rather than
    # letting python-pptx allocate fresh ids.
    for rId, rel in source.part.rels.items():
        if rId in dest.part.rels:
            continue
        dest.part.rels._rels[rId] = _Relationship(
            dest.part.rels._base_uri,
            rId,
            rel.reltype,
            target_mode=RTM.EXTERNAL if rel.is_external else RTM.INTERNAL,
            target=rel.target_ref if rel.is_external else rel.target_part,
        )

    return dest


def move_slide(prs: Presentation, from_index: int, to_index: int) -> None:
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[from_index])
    sldIdLst.insert(to_index, ids[from_index])


def drop_shape(slide, needle: str) -> None:
    shape = shape_by_text(slide, needle)
    if shape is not None:
        shape._element.getparent().remove(shape._element)


# ---------------------------------------------------------------------------
# Content
# ---------------------------------------------------------------------------

VERSION_LINE = "Version 3.0  ·  Final Delivery  ·  August 2026"


def rewrite_existing(prs: Presentation) -> None:
    s = prs.slides

    # 1 -- Title. The three pillars now name what the platform actually does.
    edit(s[0], "Version 2.0", VERSION_LINE)
    edit(s[0], "CSCD - GROUP 2", "CSCD - GROUP 2 - CAPSTONE PROJECT  · FINAL PRESENTATION")
    edit(s[0], "Multi-Vector OSINT Collection", "14 Integrated OSINT Tools, One Pipeline")
    edit(s[0], "Graph-Based Identity Correlation", "Graph Correlation with Scored Confidence")
    edit(s[0], "Chain-of-Custody Reporting", "Verifiable Evidence, Redactable Reports")

    # 2 -- Agenda, matching the sections this deck now has.
    edit(s[1], "Summary & Business Impact", "Outcome & Business Impact")
    edit(s[1], "System Architecture & Topology", "Architecture & Technology")
    edit(s[1], "Data Flow & Trust Boundaries", "Collection: Tools, Parsers, Evidence")
    edit(s[1], "Implementation Blueprint & Milestones", "Correlation, Reporting & Disclosure")
    edit(s[1], "Production-Ready Quality Standards", "Resilience, Quality & Compliance")

    # 4 -- Problem and solution, in the past tense of something delivered.
    edit(s[3], "Threat actors construct", "\n".join([
        "Threat actors compartmentalise their footprint so that no",
        "single lookup identifies them:",
        "",
        "  •  Burner / VoIP numbers, no identity document required",
        "  •  Disposable and forwarding email addresses",
        "  •  A username reused across dozens of platforms",
        "  •  Stolen or stock profile photographs",
        "",
        "Correlating this by hand costs hours per subject, and the",
        "result is a folder of tool output that no reviewer can",
        "audit: no provenance, no confidence, nothing repeatable.",
    ]))
    shrink(s[3], "Threat actors compartmentalise", 9)
    edit(s[3], "Ghost Identity Hunter automates", "\n".join([
        "Delivered: one command runs the collection, correlation",
        "and reporting end to end.",
        "",
        "  ✓  14 external OSINT tools plus 5 native modules,",
        "      each read by a parser written for its own format",
        "  ✓  Breadth-first expansion across 7 seed types with",
        "      exact-value matching, so no near miss is recorded",
        "  ✓  Graph correlation into persona clusters, every score",
        "      shown with the evidence that produced it",
        "  ✓  Every tool run preserved and SHA-256 verifiable",
        "  ✓  One HTML report: findings first, detail collapsed,",
        "      plus a masked copy that is safe to circulate",
        "  ✓  Runs entirely on the analyst's host. No cloud cost.",
    ]))
    shrink(s[3], "Delivered: one command", 9)

    # 5 -- The pipeline gained a plugin/parser layer and an evidence store.
    edit(s[4], "4-Layer Pipeline Architecture", "4-Layer Pipeline, As Built")
    edit(s[4], "CLI (Click) → Seed Validation",
         "CLI (Click) · 7 seed types · exact-match policy · BFS queue")
    edit(s[4], "Phone · Email · Username · Image · Breach Modules",
         "5 native modules + 14 tool integrations → per-tool parsers → evidence store")
    edit(s[4], "NetworkX Identity Graph → Connected Components → Confidence Scorer",
         "NetworkX graph → persona clusters → confidence scorer with provenance")
    edit(s[4], "SQLite Persistence → HTML / PDF Report → STIX 2.1 CTI Export",
         "SQLite → HTML report (+ masked copy) · JSON · CSV · PDF · pyvis graph")

    # 9 -- Components, corrected: phoneinfoga and h8mail were never integrated.
    edit(s[8], "Python deque · MAX_DEPTH=2",
         "deque BFS · configurable depth (default 2) · bounded worker pool")
    edit(s[8], "phoneinfoga (subprocess) · phonenumbers · E.164 · VoIP detection",
         "phonenumbers · E.164 validation · carrier, region and line type")
    edit(s[8], "Holehe (100+ services) · h8mail · HIBP API",
         "holehe (120+ services) · HIBP · LeakOSINT (opt-in, keyed)")
    edit(s[8], "Sherlock / Maigret · httpx · 400+ platforms · false-positive suppression",
         "sherlock · maigret · usufy · native sweep · exact-username filtering")
    edit(s[8], "Pillow (EXIF) · face_recognition (dlib)",
         "exiftool + Pillow EXIF/GPS · optional dlib face similarity")
    edit(s[8], "NetworkX DiGraph · Connected components · pyvis interactive visualization",
         "NetworkX DiGraph · weakly connected components · pyvis graph")
    edit(s[8], "Typed link weights · staleness decay · cross-type density metric",
         "Typed link weights · staleness decay · per-score provenance breakdown")
    edit(s[8], "Jinja2 HTML · WeasyPrint PDF · python-stix2 STIX 2.1 bundle",
         "Jinja2 HTML · JSON · CSV · PDF via pandoc · redacted twin file")
    edit(s[8], "Platforms probed", "Platforms probed")
    edit(s[8], "OSINT modules", "Native OSINT modules")
    edit(s[8], "Email services (Holehe)", "Email services (holehe)")
    edit(s[8], "STIX 2.1", "14")
    edit(s[8], "CTI output format", "External tools integrated")

    # 10 -- Trust boundaries: name the tools that actually cross them.
    edit(s[9], "Holehe HTTP Probes", "holehe / maigret probes")
    edit(s[9], "Sherlock / Maigret Probes", "sherlock / usufy probes")
    edit(s[9], "Reverse Image APIs", "Wayback CDX · Shodan · search engines")
    edit(s[9], "Module Dispatcher", "Tool dispatch + parsers")
    edit(s[9], "Audit Logger", "Evidence store (SHA-256)")
    edit(s[9], "TLS verify=True enforced",
         "TLS verified  ·  keys in env vars only  ·  argv lists, never a shell  "
         "·  each tool in its own process group")

    # 11 -- Technology matrix, reduced to what is imported.
    edit(s[10], "Validated Technology Matrix", "Technology Stack, As Shipped")
    edit(s[10], "phonenumbers / phoneinfoga", "phonenumbers")
    edit(s[10], "holehe + h8mail", "holehe + HIBP")
    edit(s[10], "100+ service registration detection", "120+ services; CSV output parsed")
    edit(s[10], "Sherlock / Maigret", "sherlock / maigret / usufy")
    edit(s[10], "400+ platform parallel probing", "Native NDJSON + tree output parsed")
    edit(s[10], "Pillow + face_recog", "Pillow + exiftool")
    edit(s[10], "EXIF extraction; local face similarity",
         "EXIF/GPS; dlib face match is optional")
    edit(s[10], "HIBP API v3", "HIBP + LeakOSINT")
    edit(s[10], "k-anonymity; no full hash transmitted",
         "Keyed and opt-in; token never stored in the report")
    edit(s[10], "Jinja2 + WeasyPrint + python-stix2", "Jinja2 + pandoc")
    edit(s[10], "HTML / PDF / STIX 2.1", "HTML · JSON · CSV · PDF")
    edit(s[10], "Zero-dependency; investigator-local",
         "stdlib; investigator-local; evidence blobs on disk")
    edit(s[10], "≥85% coverage; zero lint violations", "588 tests · ruff clean")

    # 14 -- Controls, described by the mechanism that exists.
    edit(s[13], "ISO 8601 timestamps; module name; SHA-256 hashed artifact values",
         "Per-run row: tool, argv, target, duration, exit status, ISO 8601 time")
    edit(s[13], "source_module + raw_evidence_ref + discovered_at per artifact row",
         "Every finding cites the run it came from; report links to that output")
    edit(s[13], "SHA-256 hash of SQLite DB computed on close; verified on report gen",
         "Each captured output hashed on write; `ghost-hunter evidence` re-verifies")
    edit(s[13], "Exp. backoff with jitter on HTTP 429; max 3 retries per endpoint",
         "Per-tool timeout, process-group kill, bounded output; a run cannot hang")
    edit(s[13], "Signed timestamped ZIP: SQLite DB + raw evidence + report on close",
         "Report + masked twin + CSV/JSON exports; raw captures kept beside the DB")

    # 15 -- The roadmap is history now.
    edit(s[14], "SECTION 4 — IMPLEMENTATION BLUEPRINT", "SECTION 4 — WHAT WAS DELIVERED")
    edit(s[14], "Delivery Milestones — 4-Week Roadmap", "Delivered Increments")
    edit(s[14], "Core Infrastructure", "Core Infrastructure")
    edit(s[14], "SQLite schema · BFS orchestrator · CLI entry point",
         "SQLite schema · BFS orchestrator · CLI · config-driven tool table")
    edit(s[14], "Phone · Email · Username · Image · Breach modules + unit tests",
         "5 native modules · 14 tool integrations · a parser per output format")
    edit(s[14], "NetworkX graph · Connected components · Confidence scorer",
         "NetworkX graph · persona clusters · scorer with provenance breakdown")
    edit(s[14], "Reporting & STIX Export", "Reporting & Disclosure")
    edit(s[14], "Jinja2 HTML · WeasyPrint PDF · python-stix2 CTI bundle",
         "Findings-first HTML · geo map · timeline · redacted twin · JSON/CSV/PDF")
    edit(s[14], "Integration & Hardening", "Evidence & Hardening")
    edit(s[14], "End-to-end tests · Injection hardening · Rate-limit compliance",
         "Evidence store · process-group isolation · 588 tests · live runs")
    for old, new in (("Wk 1–2", "Done"), ("Wk 2–3", "Done"), ("Wk 3–4", "Done")):
        while (shape := shape_by_text(s[14], old)) is not None:
            set_text(shape, new)

    # 16 -- Schema: the columns database.py actually creates. The fourth panel
    # becomes `evidence`, which is the table the deck never had.
    edit(s[15], "SECTION 4 — PHASE 1: STORAGE LAYER", "SECTION 4 — STORAGE LAYER")
    edit(s[15], "status  TEXT  CHECK(...)", "description  TEXT")
    edit(s[15], "created_at  DATETIME", "status  TEXT")
    edit(s[15], "completed_at  DATETIME", "created_at  TEXT")
    edit(s[15], "db_sha256  TEXT", "")
    edit(s[15], "artifact_type  TEXT  CHECK", "artifact_type  TEXT")
    edit(s[15], "source_module  TEXT", "source  TEXT")
    # Two panels carry a raw_evidence_ref row; the artifacts one comes first.
    edit(s[15], "raw_evidence_ref  TEXT", "metadata  TEXT (JSON)")
    edit(s[15], "platform_presence", "evidence")
    edit(s[15], "presence_id  TEXT  PK", "evidence_id  TEXT  PK")
    edit(s[15], "platform_name  TEXT", "tool / operation  TEXT")
    edit(s[15], "profile_url  TEXT", "command  TEXT")
    edit(s[15], "username  TEXT", "sha256 / byte_size  TEXT")
    edit(s[15], "is_verified  BOOLEAN", "exit_status  TEXT")
    edit(s[15], "raw_evidence_ref  TEXT", "stored_path  TEXT")
    edit(s[15], "PK = Primary Key",
         "PK = Primary Key   FK = Foreign Key   ·  also: platform_presence, "
         "audit_trail, investigation_metadata, geocode_cache")

    # 18 -- Reporting: there is no STIX exporter; say what ships.
    edit(s[17], "SECTION 4 — PHASE 4: REPORTING & CTI EXPORT",
         "SECTION 4 — REPORTING & DISCLOSURE")
    edit(s[17], "Investigation Reports & STIX 2.1 Export", "What the Report Contains")
    edit(s[17], "Executive Summary", "OSINT Overview")
    edit(s[17], "Seed artifacts; persona count; confidence",
         "KPIs, headline findings, charts, graph and map — open on load")
    edit(s[17], "Consolidated emails, phones, usernames, images",
         "Correlated personas, each score shown with its provenance")
    edit(s[17], "Platform Presence Matrix", "Platform Presence")
    edit(s[17], "Profile URLs, creation dates, verification status",
         "Profile URLs and account detail, collapsed by default")
    edit(s[17], "Interactive pyvis force-directed graph",
         "pyvis graph, in the overview and standalone")
    edit(s[17], "Breach Exposure", "Breach Records")
    edit(s[17], "Breaches per artifact; data types; pwd exposure",
         "LeakOSINT / HIBP records, reduced to the database when masked")
    edit(s[17], "Evidence Chain", "Evidence Chains")
    edit(s[17], "Ordered audit trail with raw evidence references",
         "Seed → finding path, one line each, linked to the raw output")
    edit(s[17], "STIX 2.1 Bundle", "Masked twin")
    edit(s[17], "JSON file", "HTML file")
    edit(s[17], "ThreatActor + Indicators + Relationships + Notes",
         "Values stripped at generation; safe to circulate")
    edit(s[17], "STIX 2.1 OBJECTS", "ALSO EXPORTED")
    edit(s[17], "ThreatActor SDO", "JSON")
    edit(s[17], "One per identified persona cluster",
         "Whole investigation, for SIEM / further analysis")
    edit(s[17], "Indicator SDO", "CSV")
    edit(s[17], "Per email, phone, username artifact", "Artifacts and platform presences")
    edit(s[17], "Relationship SRO", "PDF")
    edit(s[17], "indicator → indicates → threat-actor",
         "Via pandoc; sections expanded so nothing is hidden")
    edit(s[17], "Note SDO", "Graph")
    edit(s[17], "Risk indicators as structured notes", "Standalone interactive pyvis HTML")
    edit(s[17], "Bundle", "Comparison")
    edit(s[17], "All objects in STIX 2.1 envelope",
         "`--compare INV-x`: what changed since an earlier run")

    # 19 -- Quality: measured, not aspirational.
    edit(s[18], "mypy --strict enforced in CI", "Annotated public APIs; ruff-enforced")
    edit(s[18], "Zero ruff violations; pre-commit hook", "ruff clean on changed modules")
    edit(s[18], "Test Coverage", "Test Suite")
    edit(s[18], "≥85% line coverage (pytest --cov)", "588 pytest tests, all passing")
    edit(s[18], "requirements.lock via pip-compile",
         "Minimum-version pins; heavy extras optional")
    edit(s[18], "git-secrets pre-commit; no keys in VCS",
         "Keys from env only; tokens never written to a report")
    edit(s[18], "MAX 2 hops", "Default 2 hops")
    edit(s[18], "≥1,500 ms", "Rate limited")
    edit(s[18], "30 sec timeout", "Per-tool budget")
    edit(s[18], "10 sec timeout", "10 sec timeout")
    edit(s[18], "500 MB soft limit", "32 MB output cap")
    edit(s[18], "SHA-256 hashed", "Masked in reports")

    # 20 -- Closing numbers that can be defended.
    edit(s[19], "►  Automated BFS OSINT collection across 4 artifact types",
         "►  BFS collection across 7 seed types, 14 tools, 5 native modules")
    edit(s[19], "►  Graph-theoretic persona correlation with quantified confidence",
         "►  Persona correlation with every confidence score explained")
    edit(s[19], "►  STRIDE-modeled security controls mapped to NIST SP 800-53",
         "►  Every finding traceable to a hashed, re-verifiable tool run")
    edit(s[19], "►  7 MITRE ATT&CK techniques detected (T1585–T1650)",
         "►  Findings-first report, masked twin, embedded location map")
    edit(s[19], "►  Legally defensible HTML / PDF / STIX 2.1 output",
         "►  No tool can crash, hang or silence an investigation")
    edit(s[19], "►  Zero infrastructure cost · Passive OSINT only",
         "►  Passive OSINT only · runs on the analyst's host · zero cost")
    edit(s[19], "400+", "588")
    edit(s[19], "Platforms enumerated", "Tests passing")
    edit(s[19], "ATT&CK techniques covered", "External tools integrated")
    edit(s[19], "NIST controls implemented", "Dedicated output parsers")
    edit(s[19], "Infrastructur e cost", "Infrastructure cost")
    edit(s[19], "Version 2.0", VERSION_LINE)

    # 13 -- ATT&CK: the technique mapping stands, the tool names had drifted.
    edit(s[12], "Sherlock/Maigret parallel enumeration across 400+ platforms",
         "sherlock · maigret · usufy · native platform sweep, exact-match filtered")
    edit(s[12], "phoneinfoga carrier & line-type detection; VoIP flag",
         "phonenumbers carrier, region and line-type detection")
    edit(s[12], "Holehe registration detection; disposable domain classification",
         "holehe registration detection across 120+ services; disposable domains")

    # The 7/ATT&CK and 8/NIST big numbers become the tool and parser counts.
    for shape in s[19].shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == "7":
            set_text(shape, "14")
        elif shape.has_text_frame and shape.text_frame.text.strip() == "8":
            set_text(shape, "17")


# Each new slide: (clone_source_index, insert_after_index, edits)
# Sources are chosen for shape count, so the replacement text lands in a box
# the design already sized for that much content.


def add_new_slides(prs: Presentation) -> None:
    """Append the capabilities the original deck could not describe."""
    n = len(prs.slides._sldIdLst)

    # --- Tool integration and parsers -------------------------------------
    slide = clone_slide(prs, 3)
    edit(slide, "SECTION 1 — SUMMARY", "SECTION 3 — COLLECTION")
    edit(slide, "Business Impact & Objectives", "One Argv Per Tool, One Parser Per Format")
    edit(slide, "THE PROBLEM", "HOW EACH TOOL IS RUN")
    edit(slide, "Threat actors compartmentalise", "\n".join([
        "Every tool is built as an argv list -- never a shell string --",
        "in one module, so the command in the report is the command",
        "that ran:",
        "",
        "  sherlock USER --print-found --timeout 5 --no-color --no-txt",
        "  maigret USER --top-sites 150 -J ndjson -fo DIR",
        "  holehe MAIL --only-used --no-color --no-clear -C",
        "  theHarvester -d DOMAIN -b duckduckgo -f REPORT",
        "  subfinder -d DOMAIN -silent -json -timeout 10",
        "  amass enum -passive -d DOMAIN",
        "  whatweb --color=never -a 1 --log-json=OUT DOMAIN",
        "  nmap -Pn -F -sV --version-light -oX OUT IP",
        "  shodan host IP     ·     exiftool -json FILE",
        "",
        "Documented in docs/TOOL_COMMANDS.md for cross-reference.",
    ]))
    shrink(slide, "Every tool is built as an argv list", 8.5)
    edit(slide, "THE SOLUTION", "HOW ITS OUTPUT IS READ")
    edit(slide, "Delivered: one command", "\n".join([
        "17 parsers, each written for one native format, replacing",
        "the single regex that used to read every tool:",
        "",
        "  ✓  maigret NDJSON · sherlock found-list · holehe CSV",
        "  ✓  subfinder JSON · amass / sublist3r host lists",
        "  ✓  whatweb JSON log · nmap XML · exiftool JSON",
        "  ✓  whois key-value · usufy tagged entities",
        "  ✓  Wayback CDX rows · Shodan host JSON",
        "",
        "A tool's own verdict beats its exit code: holehe exits 1",
        "with a full CSV, whois exits 1 for an unregistered name.",
        "Both are answers, and both are now kept.",
        "",
        "Plugins and the orchestrator share this one path, so a",
        "tool cannot be run two different ways.",
    ]))
    shrink(slide, "17 parsers, each written for one native format", 8.5)

    # --- Tool run status --------------------------------------------------
    slide = clone_slide(prs, 12)
    edit(slide, "SECTION 3 — THREAT MODELING", "SECTION 3 — COLLECTION")
    edit(slide, "MITRE ATT&CK Coverage", "Saying Exactly Why A Tool Was Silent")
    rows = [
        ("T1585.001", "Establish Accounts — Social Media",
         "sherlock · maigret · usufy · native platform sweep, exact-match filtered",
         "output", "produced output",
         "Findings parsed, and each one cites the run that produced it"),
        ("T1585.002", "Establish Accounts — Email Accounts",
         "holehe registration detection across 120+ services; disposable domains",
         "silent", "ran and found nothing",
         "Dispatched and answered: the subject is genuinely absent there"),
        ("T1586", "Compromise Accounts",
         "HIBP breach correlation; stolen credentials to active accounts",
         "n/a", "not dispatched — wrong seed type",
         "No artifact of a type this tool accepts existed in the run"),
        ("T1598", "Phishing for Information",
         "Phone/email artifacts in social engineering from breach data",
         "off", "not dispatched — disabled",
         "Switched off in config.yaml; honoured on every dispatch path"),
        ("T1650", "Acquire — Prepaid/VoIP Numbers",
         "phonenumbers carrier, region and line-type detection",
         "timeout", "ran but timed out",
         "Killed at its own budget, together with its whole process group"),
        ("T1566.003", "Spearphishing via Service",
         "Profile image reverse search reveals stolen impersonation photos",
         "error", "ran but failed (exit n)",
         "A real failure, with the exit code shown rather than guessed at"),
        ("T1056.003", "Input Capture — Web Portal",
         "Breach data links account credentials to targeted services",
         "absent", "not installed on this host",
         "Left out of the table, so the productive count means what it says"),
    ]
    for code, name, detail, badge, label, why in rows:
        edit(slide, code, badge)
        shrink(slide, badge, 8)
        edit(slide, name, label)
        edit(slide, detail, why)

    # --- Evidence preservation -------------------------------------------
    slide = clone_slide(prs, 3)
    edit(slide, "SECTION 1 — SUMMARY", "SECTION 4 — EVIDENCE")
    edit(slide, "Business Impact & Objectives", "Every Finding Has Its Raw Output")
    edit(slide, "THE PROBLEM", "WHAT IS PRESERVED")
    edit(slide, "Threat actors compartmentalise", "\n".join([
        "For each tool run, one row and one file on disk:",
        "",
        "  tool · operation · target · exact command",
        "  tool version · captured_at · duration",
        "  exit status · sha256 · byte size · stored path",
        "",
        "The report links each finding to the run that produced",
        "it, and shows the output behind it inline.",
        "",
        "  ghost-hunter evidence --id INV-abc123",
        "",
        "re-hashes every capture and reports any file whose",
        "bytes no longer match what was recorded.",
    ]))
    shrink(slide, "For each tool run, one row and one file on disk", 9)
    edit(slide, "THE SOLUTION", "WHY IT MATTERS")
    edit(slide, "Delivered: one command", "\n".join([
        "  ✓  A reviewer can check a claim without re-running",
        "      anything -- the output is in the record",
        "",
        "  ✓  A finding that no longer reproduces is still",
        "      explainable, because the original run is kept",
        "",
        "  ✓  Tampering is detectable rather than assumed:",
        "      the hash is taken when the bytes are written",
        "",
        "  ✓  Evidence chains show the whole path from the seed",
        "      to the finding, one line per step",
        "",
        "  ✓  Citations, timeline and comparison all read the",
        "      same store, so they cannot disagree",
    ]))
    shrink(slide, "  ✓  A reviewer can check a claim", 9)

    # --- Report design ----------------------------------------------------
    slide = clone_slide(prs, 3)
    edit(slide, "SECTION 1 — SUMMARY", "SECTION 4 — REPORTING")
    edit(slide, "Business Impact & Objectives", "Lead With The Answer, Not The Evidence")
    edit(slide, "THE PROBLEM", "WHAT OPENS FIRST")
    edit(slide, "Threat actors compartmentalise", "\n".join([
        "The OSINT Overview carries everything visual:",
        "",
        "  •  KPI band: artifacts, personas, platforms, breaches,",
        "      confidence, tools that produced output",
        "  •  Headline findings, in plain sentences",
        "  •  Charts: artifact mix, confidence bands, per-tool yield",
        "  •  The relationship graph itself",
        "  •  An embedded map of every location signal",
        "",
        "Below it, 13 sections stay collapsed until asked for;",
        "only breach records open by default. Printing expands",
        "everything, so the PDF loses nothing.",
    ]))
    shrink(slide, "The OSINT Overview carries everything visual", 9)
    edit(slide, "THE SOLUTION", "DISCLOSURE CONTROL")
    edit(slide, "Delivered: one command", "\n".join([
        "Redaction removes the values at generation time, so a",
        "masked file cannot be un-masked by reading its source:",
        "",
        "  ✓  Emails, phones, names, GPS, URLs and bios masked",
        "  ✓  A leaked record is reduced to its database name",
        "  ✓  Tokens and keys never reach the document at all",
        "",
        "Every run writes both files: the working report, and",
        "<name>_redacted.html beside it for circulation.",
        "",
        "Templates: standard · executive · technical · legal.",
        "Sections can be chosen per run with --report-sections.",
    ]))
    shrink(slide, "Redaction removes the values at generation time", 9)

    # --- Resilience -------------------------------------------------------
    slide = clone_slide(prs, 11)
    edit(slide, "SECTION 3 — CYBER DEFENSE & THREAT MODELING", "SECTION 5 — RESILIENCE")
    edit(slide, "STRIDE Threat Model — Platform Threats",
         "No Tool Can End An Investigation")
    failures = [
        ("Spoofing", "Adversary lookalike domain returns false OSINT positives",
         "Runaway fork",
         "Each tool leads its own process group; the group is signalled as a whole"),
        ("Tampering", "Local filesystem access corrupts SQLite evidence chains",
         "Held-open pipe",
         "The cleanup read happens only after the group is gone, so it cannot block"),
        ("Repudiation", "No signed audit log; investigator denies query history",
         "Native abort",
         "Face work is serialised and buffer-checked; failure costs its own finding"),
        ("Info Disclosure", "PII in HTTP query params or stdout in CI/CD environments",
         "Endless output",
         "Output is capped at 32 MB and decoded leniently, so no run is lost to a byte"),
        ("Denial of Service", "Unbounded BFS traversal exhausts external API rate limits",
         "Bad output",
         "Every parser is defensive; an unreadable capture is reported, not raised"),
        ("Elev. of Privilege", "Shell metacharacters in OSINT response → subprocess inject",
         "Unforeseen",
         "The CLI reports what was found instead of a traceback; argv lists, no shell"),
    ]
    for number, (letter, old_detail, mode, control) in enumerate(failures, start=1):
        edit(slide, letter, mode)
        edit(slide, old_detail, control)
        # The single STRIDE initial has no meaning here; number the modes.
        for shape in slide.shapes:
            if (shape.has_text_frame
                    and shape.text_frame.text.strip() == "S T R I D E".split()[number - 1]):
                set_text(shape, str(number))
                break
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() in {
            "Medium", "High", "Low", "Critical"
        }:
            set_text(shape, "handled")

    # --- Exact match ------------------------------------------------------
    slide = clone_slide(prs, 3)
    edit(slide, "SECTION 1 — SUMMARY", "SECTION 3 — COLLECTION")
    edit(slide, "Business Impact & Objectives", "Full-Text Matching, Not Partial")
    edit(slide, "THE PROBLEM", "THE RULE")
    edit(slide, "Threat actors compartmentalise", "\n".join([
        "A finding is recorded only if it carries the seed's",
        "exact value. Substring and near matches are dropped,",
        "with the reason logged.",
        "",
        "  --strict-match      enforce (default)",
        "  --no-strict-match   keep partials for a wide sweep",
        "",
        "matching.allow_name_variants controls one exception:",
        "a full name may still match a reordered or initialled",
        "form. Set it to false for literal-only names.",
    ]))
    shrink(slide, "A finding is recorded only if it carries", 9)
    edit(slide, "THE SOLUTION", "WHY IT IS THE DEFAULT")
    edit(slide, "Delivered: one command", "\n".join([
        "Attribution is only as good as its weakest link. A",
        "profile that merely contains the username belongs to",
        "somebody else, and once recorded it becomes a node in",
        "the graph, a cluster member, and a confidence score.",
        "",
        "  ✓  Search engine results filtered on exact value",
        "  ✓  Tool findings filtered before persistence",
        "  ✓  Dropped candidates counted, so the sweep's reach",
        "      is visible rather than silently narrowed",
    ]))
    shrink(slide, "Attribution is only as good as its weakest link", 9)

    # --- CLI --------------------------------------------------------------
    slide = clone_slide(prs, 3)
    edit(slide, "SECTION 1 — SUMMARY", "SECTION 5 — OPERATION")
    edit(slide, "Business Impact & Objectives", "Running It")
    edit(slide, "THE PROBLEM", "COMMANDS")
    edit(slide, "Threat actors compartmentalise", "\n".join([
        "  investigate   collect from one or more seeds",
        "  correlate     re-run correlation on stored artifacts",
        "  report        render HTML / JSON / CSV / PDF",
        "  graph         standalone interactive graph",
        "  evidence      re-verify every preserved capture",
        "  list          investigations, newest first",
        "  plugins       list / info / enable / disable",
        "",
        "Seeds: --email --username --phone --full-name",
        "            --image --domain --ip",
        "",
        "Full cookbook: docs/ghost.txt",
    ]))
    shrink(slide, "  investigate   collect from one or more seeds", 9)
    edit(slide, "THE SOLUTION", "TYPICAL RUNS")
    edit(slide, "Delivered: one command", "\n".join([
        "ghost-hunter investigate -u elonmusk --use-external-tools",
        "",
        "ghost-hunter investigate -d tesla.com --depth 1 \\",
        "    --report-template technical",
        "",
        "ghost-hunter investigate -e a@b.com -p +15551234567 \\",
        "    --redact-report --report-format both",
        "",
        "ghost-hunter report --id INV-abc123 --compare INV-old",
        "",
        "ghost-hunter evidence --id INV-abc123",
        "",
        "scripts/serve_reports.sh start   # browse ./reports",
    ]))
    shrink(slide, "ghost-hunter investigate -u elonmusk", 8.5)

    print(f"  added {len(prs.slides._sldIdLst) - n} slides")


# Slide order after the additions. Cloned slides land at the end, so the
# narrative is restored here: what it does, how it is built, how it collects,
# what it keeps, how it reports, how it is defended, what shipped.
FINAL_ORDER = [
    0,   # title
    1,   # agenda
    2,   # team
    3,   # business impact
    4,   # pipeline as built
    5,   # solution architecture
    6,   # sequence diagram
    7,   # component interaction
    8,   # key components
    9,   # zero-trust data flow
    10,  # technology stack
    20,  # argv per tool, parser per format
    21,  # why a tool was silent
    25,  # full-text matching
    22,  # evidence preservation
    15,  # sqlite schema
    16,  # identity graph and confidence
    17,  # what the report contains
    23,  # findings-first layout and redaction
    11,  # STRIDE
    12,  # ATT&CK
    13,  # NIST
    24,  # resilience
    18,  # quality and constraints
    14,  # delivered increments
    26,  # running it
    19,  # summary
]


def reorder(prs: Presentation) -> None:
    sldIdLst = prs.slides._sldIdLst
    current = list(sldIdLst)
    assert sorted(FINAL_ORDER) == list(range(len(current))), "order must be a permutation"
    for element in current:
        sldIdLst.remove(element)
    for index in FINAL_ORDER:
        sldIdLst.append(current[index])


def renumber(prs: Presentation) -> None:
    """Page numbers were stale and duplicated ('2 / 16' twice, '4 / 16' four times)."""
    total = len(prs.slides._sldIdLst)
    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if "/" in text and text.replace("/", "").replace(" ", "").isdigit():
                set_text(shape, f"{index} / {total}")


def main() -> int:
    path = Path(sys.argv[1] if len(sys.argv) > 1 else "docs/Capstone.pptx")
    prs = Presentation(str(path))
    print(f"opened {path} with {len(prs.slides._sldIdLst)} slides")

    print("rewriting existing slides")
    rewrite_existing(prs)
    print("adding slides")
    add_new_slides(prs)
    print("reordering")
    reorder(prs)
    print("renumbering")
    renumber(prs)

    prs.save(str(path))
    print(f"saved {path} with {len(prs.slides._sldIdLst)} slides")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
