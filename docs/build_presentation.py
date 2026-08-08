"""
build_presentation.py
======================
Rebuilds docs/Capstone.pptx into docs/Capstone_Final.pptx:

  1. Corrects factual claims that drifted from the actual codebase
     (verified against src/, pyproject.toml, and a live pytest/ruff run).
  2. Adds four new slides that were missing from the outline:
       - Live Workflow Walkthrough
       - Performance, Testing & Results
       - Known Limitations
       - Future Roadmap
     plus a closing Q&A slide.
  3. Renumbers the "n / N" footer badges sequentially across the final deck.

All new slides are built by deep-cloning an existing on-brand slide (same
layout, same card/badge shapes, same fonts and palette pulled from the
template's own theme) and only rewriting text/severity-badge colors, so the
result inherits the master's branding exactly instead of hand-placing new
shapes.

Run:  python3 docs/build_presentation.py
"""

import copy
import re

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

SRC = "docs/Capstone.pptx"
DST = "docs/Capstone_Final.pptx"

# Palette pulled directly from the template's existing shape fills / theme.
PALETTE = {
    "teal_dark": "018090",
    "teal": "028090",
    "navy": "12263A",
    "orange": "FF6B35",
    "green": "39D353",
    "red": "FF3333",
    "white": "FFFFFF",
    "muted": "B0C4D8",
    "cyan": "00E5FF",
}

SLIDE_LAYOUT_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
)
NOTES_SLIDE_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
)
R_ATTRS = [qn("r:embed"), qn("r:link"), qn("r:id")]


# --------------------------------------------------------------------------
# Generic helpers
# --------------------------------------------------------------------------

def get_shape(slide, name):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    raise KeyError(f"shape {name!r} not found on slide")


def set_para_text(shape, text, para_index=0):
    """Overwrite one paragraph's text, keeping the first run's formatting
    (font/size/bold/color) and dropping any extra runs in that paragraph."""
    tf = shape.text_frame
    p = tf.paragraphs[para_index]
    runs = p.runs
    if not runs:
        run = p.add_run()
        run.text = text
        return
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def set_break_split_text(shape, line1, line2, para_index=0):
    """For paragraphs that hold two lines as two runs joined by a soft
    <a:br/> (rather than two separate <a:p> paragraphs) - set both runs'
    text directly, keeping the existing run formatting and the break."""
    p = shape.text_frame.paragraphs[para_index]
    runs = p.runs
    runs[0].text = line1
    runs[-1].text = line2
    for r in runs[1:-1]:
        r._r.getparent().remove(r._r)


def set_fill(shape, hex_color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex_color)


def duplicate_slide(prs, index):
    """Deep-clone prs.slides[index] and append the copy at the end of the
    deck, correctly rewriting image relationship IDs so pictures survive."""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)

    # The new slide inherits placeholder shapes from its layout - strip them
    # so we start from a blank canvas before copying the source's shapes.
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    rid_map = {}
    for rid, rel in source.part.rels.items():
        if rel.reltype in (SLIDE_LAYOUT_RELTYPE, NOTES_SLIDE_RELTYPE):
            continue
        if rel.is_external:
            new_rid = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = dest.part.rels.get_or_add(rel.reltype, rel.target_part)
        rid_map[rid] = new_rid

    for shp in source.shapes:
        new_el = copy.deepcopy(shp._element)
        for el in new_el.iter():
            for attr in R_ATTRS:
                if attr in el.attrib and el.attrib[attr] in rid_map:
                    el.attrib[attr] = rid_map[el.attrib[attr]]
        dest.shapes._spTree.append(new_el)

    return dest


def add_picture_fit(slide, image_path, left, top, max_width, max_height):
    """Place an image inside a box, preserving aspect ratio and centering it."""
    with PILImage.open(image_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    box_aspect = max_width / max_height
    if aspect > box_aspect:
        width = max_width
        height = int(max_width / aspect)
    else:
        height = max_height
        width = int(max_height * aspect)
    x = left + (max_width - width) // 2
    y = top + (max_height - height) // 2
    return slide.shapes.add_picture(image_path, x, y, width=width, height=height)


def add_caption(slide, left, top, width, text, size=10, color="12263A", bold=False):
    box = slide.shapes.add_textbox(left, top, width, 220000)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def renumber_footers(prs):
    """Find every 'n / N' page-badge text box in final slide order and
    renumber it sequentially, so newly inserted slides don't break the
    footer count (title/summary/Q&A slides have no footer, unchanged)."""
    pattern = re.compile(r"^\s*\d+\s*/\s*\d+\s*$")
    footer_shapes = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and pattern.match(shape.text_frame.text):
                footer_shapes.append(shape)
    total = len(footer_shapes)
    for i, shape in enumerate(footer_shapes, start=1):
        set_para_text(shape, f"{i} / {total}")


# --------------------------------------------------------------------------
# Load source deck
# --------------------------------------------------------------------------

prs = Presentation(SRC)
slides = prs.slides

# --------------------------------------------------------------------------
# PART 1 - Corrections to existing slides (accuracy fixes against src/)
# --------------------------------------------------------------------------

# --- Slide 1 (index 0): title slide - version number ----------------------
s = slides[0]
shp = get_shape(s, "Google Shape;134;p25")
shp.text_frame.paragraphs[0].runs[0].text = "Version 0.1.0  ·  "

# --- Slide 2 (index 1): agenda - condensed to the live 8-part run-of-show -
s = slides[1]
set_para_text(get_shape(s, "Google Shape;161;p26"), "Intro, Team & Scope of Work — Guna (2 min)")
set_para_text(get_shape(s, "Google Shape;165;p26"), "Problem Statement & OSINT Landscape — Rekha (3 min)")
set_para_text(get_shape(s, "Google Shape;169;p26"), "Architecture & Flow Walkthrough — Dhaya (5 min)")
set_para_text(get_shape(s, "Google Shape;173;p26"), "Live Tool Demos — Phone, Maigret, Rest of Tools (15 min)")
set_para_text(get_shape(s, "Google Shape;177;p26"), "Next Steps & Closing — Vindhya, Guna (2 min)")

# --- Slide 5 (index 4): 4-layer pipeline - drop fabricated STIX claim,
# tag the presenter/time for this live run-of-show section -----------------
s = slides[4]
set_para_text(
    get_shape(s, "Google Shape;255;p29"),
    "SQLite Persistence → HTML / PDF Report → JSON / CSV Export",
)
set_para_text(
    get_shape(s, "Google Shape;234;p29"),
    "SECTION 2 — ARCHITECTURAL DESIGN  ·  DHAYA  ·  5 MIN",
)

# --- Slide 4 (index 3): business impact - tag presenter/time --------------
set_para_text(
    get_shape(slides[3], "Google Shape;217;p28"),
    "SECTION 1 — SUMMARY  ·  REKHA  ·  3 MIN",
)

# --- Slide 9 (index 8): component grid -------------------------------------
s = slides[8]
set_para_text(
    get_shape(s, "Google Shape;306;p33"),
    "Holehe (100+ services) · HIBP API",
)
set_para_text(
    get_shape(s, "Google Shape;326;p33"),
    "Jinja2 HTML · Pandoc/LaTeX PDF · JSON & CSV export",
)
set_para_text(get_shape(s, "Google Shape;334;p33"), "4")
set_para_text(get_shape(s, "Google Shape;335;p33"), "Report export formats")

# --- Slide 11 (index 10): tech stack matrix --------------------------------
s = slides[10]
set_para_text(get_shape(s, "Google Shape;431;p35"), "holehe")
set_para_text(get_shape(s, "Google Shape;471;p35"), "Jinja2 + Pandoc (LaTeX)")
set_para_text(get_shape(s, "Google Shape;475;p35"), "HTML / PDF / JSON / CSV")
set_para_text(
    get_shape(s, "Google Shape;491;p35"),
    "71% coverage (572 tests); ruff cleanup in progress",
)

# --- Slide 15 (index 14): delivery roadmap - drop fabricated STIX phase --
s = slides[14]
set_para_text(get_shape(s, "Google Shape;690;p39"), "Reporting & Export")

# --- Slide 16 (index 15): DB schema - note the 4 real tables not drawn ----
s = slides[15]
legend = get_shape(s, "Google Shape;740;p40")
set_para_text(
    legend,
    "PK = Primary Key   FK = Foreign Key   CHECK = DB Constraint   "
    "·  Also present: investigation_metadata, audit_trail, evidence, geocode_cache",
)

# --- Slide 17 (index 16): correlation scoring - fix real weight ----------
s = slides[16]
set_para_text(get_shape(s, "Google Shape;773;p41"), "0.70")
# Footer box is sized for one ~60-char line (matches the original text's
# length) - leave the original staleness-decay line untouched rather than
# risk overflow; the additional link-type weights aren't a correction to
# an existing wrong value, just supplementary detail, so they're skipped.

# --- Slide 18 (index 17): reporting - remove fabricated STIX export ------
s = slides[17]
set_para_text(get_shape(s, "Google Shape;789;p42"), "SECTION 4 — PHASE 4: REPORTING & EXPORT")
set_para_text(
    get_shape(s, "Google Shape;790;p42"), "Investigation Reports & Export Formats"
)
# Report-structure table: relabel the STIX row as the real CSV export row.
set_para_text(get_shape(s, "Google Shape;822;p42"), "CSV Export")
set_para_text(get_shape(s, "Google Shape;824;p42"), "CSV file")
set_para_text(
    get_shape(s, "Google Shape;825;p42"),
    "Flat artifact & link tables for spreadsheet analysis",
)
# Right column: replace fictional "STIX 2.1 OBJECTS" block with the real
# PDF export fallback chain (pandoc -> LaTeX engines -> wkhtmltopdf -> weasyprint).
set_para_text(get_shape(s, "Google Shape;827;p42"), "PDF EXPORT PIPELINE")
set_para_text(get_shape(s, "Google Shape;828;p42"), "pandoc (primary)")
set_para_text(
    get_shape(s, "Google Shape;829;p42"), "HTML → PDF via a LaTeX engine"
)
set_para_text(get_shape(s, "Google Shape;830;p42"), "xelatex → pdflatex → lualatex")
set_para_text(get_shape(s, "Google Shape;831;p42"), "Ordered LaTeX engine fallback chain")
set_para_text(get_shape(s, "Google Shape;832;p42"), "wkhtmltopdf")
set_para_text(get_shape(s, "Google Shape;833;p42"), "Fallback if no LaTeX engine is present")
set_para_text(get_shape(s, "Google Shape;834;p42"), "weasyprint")
set_para_text(get_shape(s, "Google Shape;835;p42"), "Final fallback, pure-Python HTML→PDF")
set_para_text(get_shape(s, "Google Shape;836;p42"), "CSV / JSON")
set_para_text(get_shape(s, "Google Shape;837;p42"), "Always available, no external engine required")

# --- Slide 19 (index 18): quality standards - real, verified numbers -----
s = slides[18]
set_para_text(
    get_shape(s, "Google Shape;850;p43"), "Type hints throughout (mypy gate planned)"
)
set_para_text(
    get_shape(s, "Google Shape;852;p43"), "37 ruff findings (28 auto-fixable); cleanup in progress"
)
set_para_text(
    get_shape(s, "Google Shape;854;p43"), "71% line coverage · 572 / 572 tests passing"
)

# --- Slide 20 (index 19): closing summary ----------------------------------
s = slides[19]
set_para_text(
    get_shape(s, "Google Shape;897;p44"),
    "►  Legally defensible HTML / PDF / JSON / CSV output",
)
set_para_text(get_shape(s, "Google Shape;901;p44"), "Platforms via Sherlock/Maigret")
set_para_text(
    get_shape(s, "Google Shape;911;p44"), "Version 0.1.0  ·  Capstone Program  ·  August 2026"
)

print("Part 1 complete: corrections applied to slides 1, 5, 9, 11, 16, 17, 18, 19, 20.")

# --------------------------------------------------------------------------
# PART 2 - New slides
# --------------------------------------------------------------------------

# Donor slide indices in the *original* (pre-insertion) slide order:
ROADMAP_DONOR = 14   # Slide 15: 5-row phase-card layout
GRID_DONOR = 8        # Slide 9: 8-box component grid + 4-stat footer row
STRIDE_DONOR = 11     # Slide 12: 6-row severity-badge grid
DEMO_DONOR = 5         # Slide 6: title + single full-bleed picture

SECTION_LABEL_NAME = {
    ROADMAP_DONOR: "Google Shape;665;p39",
    GRID_DONOR: "Google Shape;292;p33",
    STRIDE_DONOR: "Google Shape;498;p36",
    DEMO_DONOR: "Google Shape;262;p30",
}
TITLE_NAME = {
    ROADMAP_DONOR: "Google Shape;666;p39",
    GRID_DONOR: "Google Shape;293;p33",
    STRIDE_DONOR: "Google Shape;499;p36",
    DEMO_DONOR: "Google Shape;263;p30",
}


def new_roadmap_slide(section_label, title, rows):
    """rows: list of 5 (tag, week, title_lines, desc) tuples, matching the
    5-phase-card layout of the original Slide 15."""
    d = duplicate_slide(prs, ROADMAP_DONOR)
    set_para_text(get_shape(d, SECTION_LABEL_NAME[ROADMAP_DONOR]), section_label)
    set_para_text(get_shape(d, TITLE_NAME[ROADMAP_DONOR]), title)

    tag_names = [
        "Google Shape;670;p39", "Google Shape;676;p39", "Google Shape;682;p39",
        "Google Shape;688;p39", "Google Shape;694;p39",
    ]
    week_names = [
        "Google Shape;671;p39", "Google Shape;677;p39", "Google Shape;683;p39",
        "Google Shape;689;p39", "Google Shape;695;p39",
    ]
    title_names = [
        "Google Shape;672;p39", "Google Shape;678;p39", "Google Shape;684;p39",
        "Google Shape;690;p39", "Google Shape;696;p39",
    ]
    desc_names = [
        "Google Shape;673;p39", "Google Shape;679;p39", "Google Shape;685;p39",
        "Google Shape;691;p39", "Google Shape;697;p39",
    ]

    for i, (tag, week, title_lines, desc) in enumerate(rows):
        set_para_text(get_shape(d, tag_names[i]), tag)
        week_shape = get_shape(d, week_names[i])
        # Collapse any multi-run week text (e.g. "Wk 1–2") to one run.
        set_para_text(week_shape, week)

        title_shape = get_shape(d, title_names[i])
        tf = title_shape.text_frame
        if len(title_lines) == 1:
            set_para_text(title_shape, title_lines[0], 0)
            # remove the now-unused second paragraph if the donor had one
            if len(tf.paragraphs) > 1:
                tf.paragraphs[1]._p.getparent().remove(tf.paragraphs[1]._p)
        else:
            set_para_text(title_shape, title_lines[0], 0)
            if len(tf.paragraphs) > 1:
                set_para_text(title_shape, title_lines[1], 1)

        set_para_text(get_shape(d, desc_names[i]), desc)

    return d


def new_grid_slide(section_label, title, cards, stats):
    """cards: list of 8 (heading, detail) tuples for the 2x4 grid.
    stats: list of 4 (big_number, label) tuples for the footer stat row."""
    d = duplicate_slide(prs, GRID_DONOR)
    set_para_text(get_shape(d, SECTION_LABEL_NAME[GRID_DONOR]), section_label)
    set_para_text(get_shape(d, TITLE_NAME[GRID_DONOR]), title)

    heading_names = [
        "Google Shape;297;p33", "Google Shape;301;p33", "Google Shape;305;p33",
        "Google Shape;309;p33", "Google Shape;313;p33", "Google Shape;317;p33",
        "Google Shape;321;p33", "Google Shape;325;p33",
    ]
    detail_names = [
        "Google Shape;298;p33", "Google Shape;302;p33", "Google Shape;306;p33",
        "Google Shape;310;p33", "Google Shape;314;p33", "Google Shape;318;p33",
        "Google Shape;322;p33", "Google Shape;326;p33",
    ]
    for i, (heading, detail) in enumerate(cards):
        set_para_text(get_shape(d, heading_names[i]), heading)
        set_para_text(get_shape(d, detail_names[i]), detail)

    stat_num_names = [
        "Google Shape;328;p33", "Google Shape;330;p33", "Google Shape;332;p33", "Google Shape;334;p33",
    ]
    stat_label_names = [
        "Google Shape;329;p33", "Google Shape;331;p33", "Google Shape;333;p33", "Google Shape;335;p33",
    ]
    for i, (num, label) in enumerate(stats):
        set_para_text(get_shape(d, stat_num_names[i]), num)
        set_para_text(get_shape(d, stat_label_names[i]), label)

    return d


def new_severity_grid_slide(section_label, title, rows):
    """rows: list of 6 (tag, heading, desc, severity) tuples, matching the
    6-row STRIDE severity-badge layout of the original Slide 12."""
    d = duplicate_slide(prs, STRIDE_DONOR)
    set_para_text(get_shape(d, SECTION_LABEL_NAME[STRIDE_DONOR]), section_label)
    set_para_text(get_shape(d, TITLE_NAME[STRIDE_DONOR]), title)

    letter_badge_names = [
        "Google Shape;502;p36", "Google Shape;509;p36", "Google Shape;516;p36",
        "Google Shape;523;p36", "Google Shape;530;p36", "Google Shape;537;p36",
    ]
    letter_text_names = [
        "Google Shape;503;p36", "Google Shape;510;p36", "Google Shape;517;p36",
        "Google Shape;524;p36", "Google Shape;531;p36", "Google Shape;538;p36",
    ]
    heading_names = [
        "Google Shape;504;p36", "Google Shape;511;p36", "Google Shape;518;p36",
        "Google Shape;525;p36", "Google Shape;532;p36", "Google Shape;539;p36",
    ]
    desc_names = [
        "Google Shape;505;p36", "Google Shape;512;p36", "Google Shape;519;p36",
        "Google Shape;526;p36", "Google Shape;533;p36", "Google Shape;540;p36",
    ]
    sev_badge_names = [
        "Google Shape;506;p36", "Google Shape;513;p36", "Google Shape;520;p36",
        "Google Shape;527;p36", "Google Shape;534;p36", "Google Shape;541;p36",
    ]
    sev_text_names = [
        "Google Shape;507;p36", "Google Shape;514;p36", "Google Shape;521;p36",
        "Google Shape;528;p36", "Google Shape;535;p36", "Google Shape;542;p36",
    ]

    sev_color = {
        "Low": PALETTE["green"],
        "Medium": PALETTE["teal"],
        "High": PALETTE["orange"],
        "Critical": PALETTE["red"],
    }

    for i, (tag, heading, desc, severity) in enumerate(rows):
        set_para_text(get_shape(d, letter_text_names[i]), tag)
        set_para_text(get_shape(d, heading_names[i]), heading)
        set_para_text(get_shape(d, desc_names[i]), desc)
        set_para_text(get_shape(d, sev_text_names[i]), severity)
        color = sev_color[severity]
        set_fill(get_shape(d, letter_badge_names[i]), color)
        set_fill(get_shape(d, sev_badge_names[i]), color)

    return d


def new_demo_slide(section_label, title, image_specs):
    """image_specs: list of (image_path, caption) shown side-by-side, real
    screenshots of the project's own generated reports. Clones the
    single-big-picture donor (Solution Architecture slide) and replaces its
    picture with 1-3 real screenshots plus captions."""
    d = duplicate_slide(prs, DEMO_DONOR)
    set_para_text(get_shape(d, SECTION_LABEL_NAME[DEMO_DONOR]), section_label)
    set_para_text(get_shape(d, TITLE_NAME[DEMO_DONOR]), title)

    big_pic = get_shape(d, "Google Shape;267;p30")
    big_pic._element.getparent().remove(big_pic._element)

    n = len(image_specs)
    margin = 342986
    top = 1090000
    bottom = 4790000
    gap = 160000
    avail_width = 9144000 - 2 * margin
    slot_width = (avail_width - gap * (n - 1)) // n
    slot_height = bottom - top - 260000

    for i, (path, caption) in enumerate(image_specs):
        slot_left = margin + i * (slot_width + gap)
        add_picture_fit(d, path, slot_left, top, slot_width, slot_height)
        add_caption(d, slot_left, top + slot_height + 30000, slot_width, caption)

    return d


def new_tools_appendix_slide(section_label, title, cards, footnote):
    """cards: up to 8 (heading, detail) tuples for the 2x4 'how it works'
    grid; unused slots are removed. The 4-stat footer band is repurposed as
    a single full-width footnote (real screenshots live on the companion
    'reference screenshots' appendix slide - shrinking them to fit this
    footer band would make the report tables illegible)."""
    d = duplicate_slide(prs, GRID_DONOR)
    set_para_text(get_shape(d, SECTION_LABEL_NAME[GRID_DONOR]), section_label)
    set_para_text(get_shape(d, TITLE_NAME[GRID_DONOR]), title)

    bg_names = [
        "Google Shape;295;p33", "Google Shape;299;p33", "Google Shape;303;p33",
        "Google Shape;307;p33", "Google Shape;311;p33", "Google Shape;315;p33",
        "Google Shape;319;p33", "Google Shape;323;p33",
    ]
    accent_names = [
        "Google Shape;296;p33", "Google Shape;300;p33", "Google Shape;304;p33",
        "Google Shape;308;p33", "Google Shape;312;p33", "Google Shape;316;p33",
        "Google Shape;320;p33", "Google Shape;324;p33",
    ]
    heading_names = [
        "Google Shape;297;p33", "Google Shape;301;p33", "Google Shape;305;p33",
        "Google Shape;309;p33", "Google Shape;313;p33", "Google Shape;317;p33",
        "Google Shape;321;p33", "Google Shape;325;p33",
    ]
    detail_names = [
        "Google Shape;298;p33", "Google Shape;302;p33", "Google Shape;306;p33",
        "Google Shape;310;p33", "Google Shape;314;p33", "Google Shape;318;p33",
        "Google Shape;322;p33", "Google Shape;326;p33",
    ]
    for i in range(8):
        if i < len(cards):
            set_para_text(get_shape(d, heading_names[i]), cards[i][0])
            set_para_text(get_shape(d, detail_names[i]), cards[i][1])
        else:
            for nm in (bg_names[i], accent_names[i], heading_names[i], detail_names[i]):
                shp = get_shape(d, nm)
                shp._element.getparent().remove(shp._element)

    stat_names = [
        "Google Shape;328;p33", "Google Shape;329;p33", "Google Shape;330;p33",
        "Google Shape;331;p33", "Google Shape;332;p33", "Google Shape;333;p33",
        "Google Shape;334;p33", "Google Shape;335;p33",
    ]
    for nm in stat_names[1:]:
        shp = get_shape(d, nm)
        shp._element.getparent().remove(shp._element)
    footnote_shape = get_shape(d, stat_names[0])
    footnote_shape.left = get_shape(d, "Google Shape;327;p33").left
    footnote_shape.top = get_shape(d, "Google Shape;327;p33").top
    footnote_shape.width = get_shape(d, "Google Shape;327;p33").width
    set_para_text(footnote_shape, footnote)
    for run in footnote_shape.text_frame.paragraphs[0].runs:
        run.font.size = Pt(11)
        run.font.bold = False

    return d


# --- New slide: Confidence Scoring Algorithm — How It Works ---------------
# Pulled directly from src/correlation/scorer.py (compute_link_confidence,
# explain_identity_risk_score, classify_risk_level) - note the risk-level
# thresholds here (0.8/0.6/0.4/0.2) are the real ones from the function
# body; the module's own docstring states different, stale thresholds
# (0.9/0.7/0.5), which is exactly the kind of drift this deck now avoids.
scoring_slide = new_roadmap_slide(
    "SECTION 4 — PHASE 3: CORRELATION ENGINE",
    "Confidence Scoring Algorithm — How It Works",
    [
        (
            "1", "Base",
            ["Base Score by Link Type"],
            "BASE_SCORES lookup: exact_match 1.00 down to temporal_match 0.40 (scorer.py)",
        ),
        (
            "2", "Decay",
            ["Freshness Decay"],
            ">365 days → ×0.80  |  >730 days → ×0.60, applied to the base score",
        ),
        (
            "3", "×Src",
            ["Source Reliability"],
            "score = base × source_reliability (0.0–1.0, defaults to 1.0)",
        ),
        (
            "4", "Clamp",
            ["Bound & Round"],
            "round(min(score, 1.0), 3) → final per-link confidence value",
        ),
        (
            "5", "ΣCap",
            ["Risk Aggregation"],
            "Σ indicator weights, capped at 1.0 → critical ≥.8, high ≥.6, medium ≥.4, low ≥.2",
        ),
    ],
)

# --- New slide: Live Workflow Walkthrough ----------------------------------
workflow_slide = new_roadmap_slide(
    "SECTION 5 — LIVE WORKFLOW WALKTHROUGH",
    "End-to-End Investigation — CLI Walkthrough",
    [
        (
            "Step 1", "Seed",
            ["Launch Investigation"],
            "ghost-hunter investigate --phone +1... --email ... --username ... --depth 2",
        ),
        (
            "Step 2", "Collect",
            ["Parallel Collection"],
            "Orchestrator dispatches phone/email/username/image plugins across the BFS frontier",
        ),
        (
            "Step 3", "Correlate",
            ["Graph Correlation"],
            "ghost-hunter correlate — builds NetworkX graph, scores links, decays stale evidence",
        ),
        (
            "Step 4", "Visualize",
            ["Graph Review"],
            "ghost-hunter graph — opens the interactive pyvis identity graph in-browser",
        ),
        (
            "Step 5", "Report",
            ["Export Evidence"],
            "ghost-hunter report --report-format html,pdf,json,csv  +  ghost-hunter evidence",
        ),
    ],
)

# --- New slide: Performance, Testing & Results ------------------------------
perf_slide = new_grid_slide(
    "SECTION 5 — VERIFIED RESULTS",
    "Performance, Testing & Results",
    [
        ("Test Suite", "572 tests · 0 failures · ~32–46s full run (pytest)"),
        ("Line Coverage", "71% overall (7,327 statements) — measured via pytest --cov"),
        ("Weakest Modules", "exports.py 32% · http_client.py 61% · tool_checker.py 76%"),
        ("Static Analysis", "ruff: 37 findings, 28 auto-fixable — cleanup in progress"),
        ("Type Checking", "Type hints throughout source; mypy CI gate not yet wired up"),
        (
            "Perf Optimization",
            "Parallel BFS cut a representative run from hours to ~21s "
            "(partial run, external tools disabled)",
        ),
        ("CI Gate", "pytest + ruff on pre-commit; mypy and coverage gate planned"),
        ("Investigation Depth", "MAX_DEPTH=2 default · 30s subprocess / 10s HTTP timeouts"),
    ],
    [
        ("572", "Tests passing"),
        ("71%", "Line coverage"),
        ("37", "Open ruff findings"),
        ("~21s", "Runtime (optimized, partial run)"),
    ],
)

# --- New slide: Known Limitations -------------------------------------------
limitations_slide = new_severity_grid_slide(
    "SECTION 6 — HONEST ASSESSMENT",
    "Known Limitations",
    [
        (
            "1", "Provider Coverage Gap",
            "4 keyed API providers + 22 local CLI plugins vs. Maltego's 99+ transforms",
            "High",
        ),
        (
            "2", "No CTI/STIX Export",
            "Threat-actor / indicator export (STIX 2.1) is designed but not yet implemented",
            "Medium",
        ),
        (
            "3", "Concurrency Risks",
            "Nested thread pools can oversubscribe; rate-limiter holds its lock during sleep",
            "High",
        ),
        (
            "4", "No Type-Check Gate",
            "mypy is not wired into CI yet, despite type hints being present throughout",
            "Medium",
        ),
        (
            "5", "Under-tested Modules",
            "exports.py (32%) and http_client.py (61%) fall well below the 71% average",
            "Medium",
        ),
        (
            "6", "Face-Match Perf",
            "CPU-bound face_recognition runs on threads, not a ProcessPoolExecutor",
            "Low",
        ),
    ],
)

# --- New slide: Future Roadmap ----------------------------------------------
roadmap_slide = new_roadmap_slide(
    "SECTION 6 — WHAT'S NEXT",
    "Future Roadmap — Next Phases",
    [
        (
            "Next", "P1",
            ["Concurrency", "Hardening"],
            "Fix rate-limiter lock-during-sleep bug; move Google Dorks onto the pooled HTTP session",
        ),
        (
            "Next", "P2",
            ["Performance"],
            "ProcessPoolExecutor for CPU-bound face matching; async I/O for platform fan-out",
        ),
        (
            "Next", "P3",
            ["CTI Export"],
            "Implement STIX 2.1 bundle export — ThreatActor / Indicator / Relationship / Note SDOs",
        ),
        (
            "Next", "P4",
            ["Quality Gate"],
            "Add mypy --strict to CI; raise coverage on exports.py and http_client.py toward 85%",
        ),
        (
            "Next", "P5",
            ["Provider Growth"],
            "Expand keyed-API and plugin coverage toward parity with commercial OSINT platforms",
        ),
    ],
)

# --- New slide: Q&A closer (cloned from the title slide's layout) ---------
qa_slide = duplicate_slide(prs, 0)
set_para_text(get_shape(qa_slide, "Google Shape;130;p25"), "CSCD — GROUP 2 · CAPSTONE PROJECT", 0)
title_shape = get_shape(qa_slide, "Google Shape;131;p25")
set_break_split_text(title_shape, "Questions", "& Discussion")
set_para_text(get_shape(qa_slide, "Google Shape;132;p25"), "Ghost Identity Hunter — Thank You", 0)
if len(get_shape(qa_slide, "Google Shape;132;p25").text_frame.paragraphs) > 1:
    p = get_shape(qa_slide, "Google Shape;132;p25").text_frame.paragraphs[1]
    p._p.getparent().remove(p._p)
set_para_text(get_shape(qa_slide, "Google Shape;134;p25"), "Version 0.1.0  ·  Capstone Program", 0)
set_para_text(get_shape(qa_slide, "Google Shape;139;p25"), "Full documentation in /docs")
set_para_text(get_shape(qa_slide, "Google Shape;144;p25"), "Live CLI demo available on request")
set_para_text(get_shape(qa_slide, "Google Shape;149;p25"), "Open floor for questions & discussion")

# --- New slides: sequence walkthrough — what talks to what, end to end ----
# Sourced from docs/SEQUENCE_DIAGRAMS.md (real mermaid sequenceDiagrams that
# mirror src/orchestrator.py and src/reporting/html_report.py).
seq1_slide = new_roadmap_slide(
    "SECTION 2 — ARCHITECTURAL DESIGN  ·  DHAYA",
    "Sequence Walkthrough — End-to-End Investigation",
    [
        (
            "1", "CLI",
            ["Seed & Dispatch"],
            "cli.py investigate calls orchestrator.run_investigation(); seeds load into the BFS queue",
        ),
        (
            "2", "BFS",
            ["Per-Artifact Processing"],
            "Each artifact hits its OSINT module, external-tool dispatch, and the plugin system in parallel",
        ),
        (
            "3", "Discover",
            ["Expand the Frontier"],
            "New discovered artifacts are deduped, persisted, and queued for the next BFS depth level",
        ),
        (
            "4", "Score",
            ["Correlate & Finalize"],
            "analyze_correlation() builds the identity graph; results are stored in investigation_metadata",
        ),
        (
            "5", "Report",
            ["Auto-Generate Report"],
            "generate_html_report() / generate_json_report() render the finished investigation to reports/",
        ),
    ],
)

seq2_slide = new_roadmap_slide(
    "SECTION 2 — ARCHITECTURAL DESIGN  ·  DHAYA",
    "Sequence Walkthrough — How the Report Gets Produced",
    [
        (
            "1", "Read",
            ["Load Investigation"],
            "cli.py report reads artifacts, links, platform presences and the audit trail from SQLite",
        ),
        (
            "2", "Link",
            ["Correlate Identities"],
            "linker.correlate_identities() rebuilds the identity graph and groups connected components",
        ),
        (
            "3", "Score",
            ["Score & Classify"],
            "compute_identity_risk_score() and classify_risk_level() attach risk to each identity profile",
        ),
        (
            "4", "Render",
            ["Build Sections & Graph"],
            "Timeline, key findings, risk matrix, heatmap, and an embedded pyvis relationship graph are built",
        ),
        (
            "5", "Write",
            ["Render & Save"],
            "Jinja2 renders the selected template; HTML is written to reports/{investigation_id}_report.html",
        ),
    ],
)

# --- New slides: live tool demos, real screenshots from actual CLI runs ---
# Every screenshot below is a genuine Chrome-headless capture of a report
# this project's own `ghost-hunter investigate` produced against safe,
# non-personal test inputs (a libphonenumber demo number, a throwaway demo
# username, and the RFC 2606 reserved test@example.com address) - not a
# mockup.
phone_demo_slide = new_demo_slide(
    "LIVE DEMO — PHONE NUMBER LOOKUP  ·  KANTI  ·  5 MIN",
    "Phone Number Lookup — Real Investigation Report",
    [
        ("docs/assets/demo_phone_overview.png", "Investigation overview & risk dashboard"),
        ("docs/assets/demo_phone_evidence.png", "Evidence basis + PhoneValidationPlugin findings"),
    ],
)

maigret_demo_slide = new_demo_slide(
    "LIVE DEMO — MAIGRET FLOW  ·  SOUMO  ·  5 MIN",
    "Maigret Flow — Real Investigation Report",
    [
        ("docs/assets/demo_maigret_overview.png", "Investigation overview — 9 artifacts, 7 platforms"),
        ("docs/assets/demo_maigret_toolmetrics.png", "Per-tool breakdown — Sherlock, Google Dorks, Maigret"),
    ],
)

rest_demo_slide = new_demo_slide(
    "LIVE DEMO — REST OF THE TOOLS  ·  VINDHYA  ·  5 MIN",
    "Holehe, WHOIS & the Rest of the Toolkit",
    [
        ("docs/assets/demo_rest_overview.png", "Investigation overview — test@example.com"),
        ("docs/assets/demo_rest_holehe.png", "Per-tool breakdown — Holehe found 15 real registered accounts"),
    ],
)

# --- Appendix divider (built straight from the master's SECTION_HEADER
# layout, not cloned from a content slide) ----------------------------------
appendix_divider = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[1])
appendix_divider.placeholders[0].text = "Appendix"
sub_box = appendix_divider.shapes.add_textbox(685972, 2200000, 7772056, 700000)
sub_tf = sub_box.text_frame
sub_tf.word_wrap = True
sub_p = sub_tf.paragraphs[0]
sub_p.alignment = PP_ALIGN.CENTER
sub_run = sub_p.add_run()
sub_run.text = (
    "Supporting technical detail — threat modeling, schema, scoring algorithm, "
    "quality metrics, limitations and tool internals"
)
sub_run.font.size = Pt(14)
sub_run.font.name = "Calibri"
sub_run.font.color.rgb = RGBColor.from_string(PALETTE["muted"])

# --- Appendix: How Each OSINT Tool Works -----------------------------------
tools_appendix_slide = new_tools_appendix_slide(
    "APPENDIX — TOOL INTERNALS",
    "How Each OSINT Tool Works",
    [
        (
            "Phone Validation",
            "libphonenumber parse · carrier/region from the number, no API call",
        ),
        (
            "Sherlock",
            "~400 site GETs · per-site 200/404 rule = found/not-found",
        ),
        (
            "Maigret",
            "Same probe technique, top 150 sites · run with --no-recursion",
        ),
        (
            "Holehe",
            "Password-reset/signup endpoint probe per site · infers account existence",
        ),
        (
            "theHarvester",
            "Passive domain OSINT via search engines + crt.sh · no direct probing",
        ),
    ],
    "Real screenshots are on the next slide — not mockups. theHarvester wasn't installable "
    "here, so it's description-only above.",
)

# --- Appendix: reference screenshots (full-size, for Q&A backup) ----------
tools_screenshots_slide = new_demo_slide(
    "APPENDIX — TOOL INTERNALS",
    "Tool Output — Reference Screenshots",
    [
        ("docs/assets/demo_phone_evidence.png", "Phone Validation — evidence & confidence"),
        ("docs/assets/demo_maigret_toolmetrics.png", "Sherlock + Maigret — per-tool breakdown"),
        ("docs/assets/demo_rest_holehe.png", "Holehe — per-tool breakdown"),
    ],
)

print("Part 2 complete: 14 new slides built (scoring, workflow, performance, "
      "limitations, roadmap, Q&A, 2 sequence walkthroughs, 3 tool demos, "
      "appendix divider, appendix tools + screenshots).")

# --------------------------------------------------------------------------
# PART 3 - Reorder into: Live 8-part run-of-show, then Appendix
# --------------------------------------------------------------------------
# Every new slide was appended to the end of the deck by duplicate_slide, in
# this creation order. Name each by its append position instead of chaining
# pop/insert moves - easier to get right and to re-check.
n_original = 20  # indices 0-19 = the original slides 1..20, unchanged
(
    scoring_i, workflow_i, perf_i, limitations_i, roadmap_i, qa_i,
    seq1_i, seq2_i, phone_demo_i, maigret_demo_i, rest_demo_i,
    appendix_divider_i, tools_appendix_i, tools_screenshots_i,
) = range(n_original, n_original + 14)

# LIVE DECK - matches the presenter run-of-show exactly:
#   Guna(2m): Title, Agenda, Team
#   Rekha(3m): Business Impact, Key Components & Tech (OSINT tools)
#   Dhaya(5m): 4-Layer Pipeline, Solution Architecture, Sequence Diagram,
#              2 NEW sequence walkthroughs, Component Interaction,
#              Zero Trust Data Flow, Tech Stack Matrix
#   Kanti(5m): Phone Number Lookup demo
#   Soumo(5m): Maigret Flow demo
#   Vindhya(5m+1m): Rest of the Tools demo, Future Roadmap (Next Steps)
#   Guna(1m): Summary, Q&A
live_deck = [
    0, 1, 2,                                    # Title, Agenda, Team
    3, 8,                                        # Business Impact, Key Components & Tech
    4, 5, 6, seq1_i, seq2_i, 7, 9, 10,            # Architecture & Flow Walkthrough
    phone_demo_i, maigret_demo_i, rest_demo_i,    # Live tool demos
    roadmap_i,                                    # Follow-up & Next Steps
    19, qa_i,                                     # Summary, Q&A
]

# APPENDIX - backup/reference material, not part of the timed live run.
appendix = [
    appendix_divider_i,
    11, 12, 13,                       # STRIDE, MITRE ATT&CK, NIST controls
    14, 15, 16, scoring_i,            # Delivery roadmap, DB schema, correlation + scoring algo
    17, 18,                           # Reporting & Export, Quality Standards
    perf_i, limitations_i,            # Performance/Testing/Results, Known Limitations
    workflow_i,                       # General CLI reference (Live Workflow Walkthrough)
    tools_appendix_i, tools_screenshots_i,  # How each tool works + reference screenshots
]

final_order = live_deck + appendix
assert sorted(final_order) == list(range(n_original + 14)), "final_order must be a full permutation"

sldIdLst = prs.slides._sldIdLst
id_elements = list(sldIdLst)
for el in id_elements:
    sldIdLst.remove(el)
for idx in final_order:
    sldIdLst.append(id_elements[idx])

print("Part 3 complete: slide order finalized (live deck + appendix).")

# --------------------------------------------------------------------------
# PART 4 - Renumber footers and save
# --------------------------------------------------------------------------
renumber_footers(prs)

prs.save(DST)
print(f"Saved {DST} with {len(prs.slides)} slides.")